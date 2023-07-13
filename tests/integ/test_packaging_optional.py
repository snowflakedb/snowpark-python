#
# Copyright (c) 2012-2023 Snowflake Computing Inc. All rights reserved.
#

from unittest.mock import patch

import pytest

from snowflake.snowpark import Row
from snowflake.snowpark.exceptions import SnowparkSQLException
from snowflake.snowpark.functions import udf
from tests.utils import TempObjectType, Utils

permanent_stage_name = "permanent_stage_for_package_testing"
reinstall_options = [True, False]
# Note: Tests will run much faster if only False is included in reinstall_options.
# Only including False will test the UDF snippets without pip installing packages.


@pytest.fixture(scope="module", autouse=True)
def setup(session, resources_path):
    try:
        print(
            f"Files on permanent stage: {session._list_files_in_stage(permanent_stage_name)}"
        )
    except SnowparkSQLException as sse:
        if "does not exist or not authorized" not in str(sse):
            raise sse
        Utils.create_stage(session, permanent_stage_name, is_temporary=False)


@pytest.fixture(autouse=True)
def clean_up(session):
    # Note: All tests in this module are skipped as these tests are only intended for in-depth testing of packaging.
    # Please run these tests when any change to packaging functionality is made.
    pytest.skip()

    session.clear_packages()
    session.clear_imports()
    yield


@pytest.mark.parametrize("force_install", reinstall_options)
def test_sktime(session, force_install):
    """
    Assert that sktime package is usable by running an K Neighbors Time Series classifier on random data.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["sktime", "pmdarima"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def sktime_test() -> str:
            import numpy as np
            from sktime.classification.distance_based import (
                KNeighborsTimeSeriesClassifier,
            )

            # Generate random time series data
            n_samples = 100  # Number of time series samples
            n_timestamps = 50  # Number of timestamps per time series
            n_classes = 3  # Number of classes
            X = np.random.rand(n_samples, n_timestamps)
            y = np.random.randint(n_classes, size=n_samples)

            # Create a K-Nearest Neighbors Time Series Classifier
            classifier = KNeighborsTimeSeriesClassifier(n_neighbors=3)

            # Fit the classifier to the data
            classifier.fit(X, y)

            # Make predictions on new data
            X_test = np.random.rand(5, n_timestamps)
            y_pred = classifier.predict(X_test)

            return str(len(y_pred))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("5")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_skfuzzy(session, force_install):
    """
    Assert that scikit-fuzzy package is usable by running c-means clustering.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["scikit-fuzzy"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def skfuzzy_test() -> str:
            import numpy as np
            import skfuzzy as fuzz

            # Generate random data
            np.random.seed(0)
            x = np.arange(0, 10, 0.1)  # Input variable
            x_len = len(x)
            y = np.random.randn(x_len)  # Output variable

            # Apply fuzzy c-means clustering
            cntr, u, _, _, _, _, _ = fuzz.cluster.cmeans(
                data=np.vstack([x, y]),
                c=2,  # Number of clusters
                m=2,  # Fuzziness coefficient
                error=0.005,
                maxiter=1000,
                init=None,
            )

            # Get the membership values of the data points for each cluster
            cluster_membership = np.argmax(u, axis=0)
            return str(len(cluster_membership))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("100")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_pyod(session, force_install):
    """
    Assert that pyod package is usable by running a KNN classifier.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["pyod"], persist_path=permanent_stage_name, force_install=force_install
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def pyod_test() -> str:
            import numpy as np
            from pyod.models.knn import KNN

            # Generate random data
            np.random.seed(0)
            X_train = np.random.randn(1000, 2)  # 1000 samples with 2 features

            # Create and fit the KNN outlier detection model
            model = KNN(
                contamination=0.1
            )  # contamination represents the expected proportion of outliers
            model.fit(X_train)

            # Generate test data with outliers
            X_test = np.random.randn(200, 2)  # 200 samples with 2 features
            outlier_scores = model.decision_function(
                X_test
            )  # Obtain outlier scores for test data

            return str(len(outlier_scores))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("200")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_parsy(session, force_install):
    """
    Assert that parsy package is usable by parsing a dd-mm-yy date.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["parsy"], persist_path=permanent_stage_name, force_install=force_install
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def parsy_test() -> str:
            from datetime import date

            from parsy import regex, string

            ddmmyy = (
                regex(r"[0-9]{2}")
                .map(int)
                .sep_by(string("-"), min=3, max=3)
                .combine(lambda d, m, y: date(2000 + y, m, d))
            )
            return str(ddmmyy.parse("06-05-14"))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("2014-05-06")],
        )


@pytest.mark.xfail(reason="atpublic is not usable due to KeyError issues with __all__")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_atpublic(session, force_install):
    """
    Assert that atpublic package is usable by making a 'public' decorated function.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["atpublic"], persist_path=permanent_stage_name, force_install=force_install
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def atpublic_test() -> str:
            from public import public

            @public
            def foo():
                pass

            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


def test_tiktoken(session):
    """
    Assert that tiktoken package is not usable as it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["tiktoken"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_ibis(session, force_install):
    """
    Assert that ibis package is usable by creating a simple 'and' template.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["ibis"], persist_path=permanent_stage_name, force_install=force_install
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def ibis_test() -> str:
            import ibis

            template = ibis.Template("{{foo}} and {{bar}}")
            return template.render({"foo": "ham", "bar": "eggs"})

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("ham and eggs")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_np_financial(session, force_install):
    """
    Assert that numpy-financial package is usable by calculating present discounted cash flow value.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["numpy-financial"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def numpy_financial_test() -> str:
            import numpy as np
            import numpy_financial as npf

            # Generate cash flows
            cash_flows = np.array([-100, 50, 40, 30, 20])

            # Set discount rate
            discount_rate = 0.1

            # Calculate present value
            present_value = npf.npv(discount_rate, cash_flows)
            return str(round(present_value, 2))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("14.71")],
        )


@pytest.mark.xfail(
    reason="monai is not usable due to psutil.__version__ attribute missing"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_monai(session, force_install):
    """
    Assert that monai package is usable by transforming images.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["monai", "pytorch", "psutil==5.9.5"],
            persist_path=permanent_stage_name,
            force_push=True,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def monai_test() -> str:
            import torch
            from monai.data import DataLoader
            from monai.transforms import Compose, ScaleIntensity

            # Create a dummy dataset
            class DummyDataset(torch.utils.data.Dataset):
                def __init__(self) -> None:
                    self.data = torch.randn((100, 1, 64, 64))
                    return

                def __getitem__(self, index):
                    return self.data[index]

                def __len__(self):
                    return len(self.data)

            # Initialize the dataset and data loader
            dataset = DummyDataset()
            dataloader = DataLoader(dataset, batch_size=8, num_workers=0)

            # Define the transform pipeline
            transform = Compose([ScaleIntensity()])

            # Apply transform to each batch in the data loader
            for batch in dataloader:
                transformed_batch = transform(batch)
                print(transformed_batch.shape)
            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_textdistance(session, force_install):
    """
    Assert that textdistance package is usable by calculating normalized levenshtein similarity.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["textdistance"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def textdistance_test() -> str:
            import textdistance

            return str(
                textdistance.levenshtein.normalized_similarity("text", "distance")
            )

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("0.125")],
        )


def test_faiss(session):
    """
    Assert that faiss package is not usable because it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["tiktoken"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


def test_duckdb(session):
    """
    Assert that duckdb package is not usable because it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["duckdb"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_asana(session, force_install):
    """
    Assert that asana package is usable by attempting an oauth connection to Asana.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["asana"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def asana_test() -> str:
            import asana

            client = asana.Client.oauth(
                client_id="ASANA_CLIENT_ID",
                client_secret="ASANA_CLIENT_SECRET",
                redirect_uri="https://yourapp.com/auth/asana/callback",
            )
            (url, state) = client.session.authorization_url()
            return str(len(state))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("30")],
        )


def test_whylogs(session):
    """
    Assert that whylogs package is not usable because it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["whylogs"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_deepdiff(session, force_install):
    """
    Assert that deepdiff package is usable by diffing between two dictionaries.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["deepdiff"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def deepdiff_test() -> str:
            from deepdiff import DeepDiff

            # Sample dictionaries
            dict1 = {"a": 1, "b": 2, "c": 3}
            dict2 = {"a": 1, "b": 2, "c": 4}

            # Perform deep comparison
            diff = DeepDiff(dict1, dict2)

            # Print the differences
            return diff

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [
                Row(
                    "{'values_changed': {\"root['c']\": {'new_value': 4, 'old_value': 3}}}"
                )
            ],
        )


def test_quantile_forest(session):
    """
    Assert that quantile_forest package is not usable because it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["quantile_forest", "setuptools"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.xfail(
    reason="pysoundfile is not usable because it needs the libsndfile C library to be present."
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_pysoundfile(session, force_install):
    """
    Assert that PySoundFile package is usable by importing it.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["PySoundFile"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def pysoundfile_test() -> str:
            import soundfile as sf

            print(sf.__version__)
            return "works"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("works")],
        )


def test_librosa(session):
    """
    Assert that librosa package is unusable because of native dependencies.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["librosa"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_sqlglot(session, force_install):
    """
    Assert that sqlglot package is usable by converting duckdb SQL into a Hive query.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["sqlglot"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def sqlglot_test() -> str:
            import sqlglot

            return sqlglot.transpile(
                "SELECT EPOCH_MS(1618088028295)", read="duckdb", write="hive"
            )[0]

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("SELECT FROM_UNIXTIME(1618088028295 / 1000)")],
        )


@pytest.mark.xfail(reason="Pip error")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_greykite(session, force_install):
    """
    Assert that greykite package is usable by performing a time series forecasting task.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["greykite", "cmake"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def greykite_test() -> str:
            from greykite.framework.templates.forecaster import Forecaster

            return str(Forecaster)

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_text2num(session, force_install):
    """
    Assert that text2num package is usable by converting Spanish number text to digits.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["text2num"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def text2num_test() -> str:
            from text_to_num import text2num

            return str(text2num("nueve mil novecientos noventa y nueve", "es"))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("9999")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_scrubadub(session, force_install):
    """
    Assert that scrubadub package is usable by cleaning PII information from text.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["scrubadub"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def scrubadub_test() -> str:
            import scrubadub

            return scrubadub.clean("Contact Vivek Nayak at vivek.nayak@snowflake.com")

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("Contact Vivek {{EMAIL}}@snowflake.com")],
        )


@pytest.mark.xfail(reason="Pip error")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_mecab(session, force_install):
    """
    Assert that mecab package is usable by parsing Japanese.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["mecab-python3"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def mecab_test() -> str:
            import MeCab

            wakati = MeCab.Tagger("-Owakati")
            return wakati.parse("pythonが大好きです").split()[2]

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("大好き")],
        )


@pytest.mark.xfail(reason="Unable to import pptx file from the zipped environment.")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_python_pptx(session, force_install):
    """
    Assert that python-pptx package is usable by creating a presentation
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["python-pptx"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def pptx_test() -> str:
            from pptx import Presentation

            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Hello, World!"
            subtitle.text = "python-pptx was here!"
            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_imodels(session, force_install):
    """
    Assert that imodels is usable by training a hierarchical sparse decision tree.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["imodels"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def imodels_test() -> str:
            import numpy as np
            from imodels import HSTreeClassifierCV
            from sklearn.model_selection import train_test_split

            # Generate synthetic data
            np.random.seed(42)
            n_samples = 1000
            n_features = 10
            X = np.random.randn(n_samples, n_features)
            y = np.random.randint(0, 2, size=n_samples)
            feature_names = [f"Feature_{i}" for i in range(n_features)]

            # Split the data into training and testing sets
            X_train, X_test, y_train, y_test = train_test_split(X, y, random_state=42)

            # Fit the model
            model = HSTreeClassifierCV(max_leaf_nodes=4)
            model.fit(X_train, y_train, feature_names=feature_names)

            # Make predictions
            preds_proba = model.predict_proba(X_test)

            return str(len(preds_proba))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("250")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_kedro(session, force_install):
    """
    Assert that kedro is usable by defining a square root node in a pipeline.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["kedro"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def kedro_test() -> str:
            from kedro.pipeline import Pipeline, node

            # Define a simple node
            def square(x):
                return x**2

            # Create a pipeline
            pipeline = Pipeline([node(square, "input", "output")])

            return str(pipeline)

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("Pipeline([\nNode(square, 'input', 'output', None)\n])")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_polars(session, force_install):
    """
    Assert that polars package is not usable as it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["polars"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_coboljsonifier(session, force_install):
    """
    Assert that coboljsonifier is usable by importing it.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["coboljsonifier"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def coboljsonifier_test() -> str:
            from coboljsonifier.config.parser_type_enum import ParseType
            from coboljsonifier.copybookextractor import CopybookExtractor
            from coboljsonifier.parser import Parser

            print(str(ParseType), str(CopybookExtractor), str(Parser))

            return "works"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("works")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_sklearn_crfsuite(session, force_install):
    """
    Assert that sklearn_crfsuite is usable by training a CRF model.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["sklearn-crfsuite", "scikit-learn"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def crfsuite_test() -> str:
            import sklearn_crfsuite

            # Generate sample data
            X_train = [
                [("feature1", "value1"), ("feature2", "value2")],
                [("feature3", "value3"), ("feature4", "value4")],
            ]
            y_train = [["label1", "label2"], ["label3", "label4"]]
            X_test = [
                [("feature5", "value5"), ("feature6", "value6")],
                [("feature7", "value7"), ("feature8", "value8")],
            ]

            # Define and train the CRF model
            crf_model = sklearn_crfsuite.CRF()
            crf_model.fit(X_train, y_train)

            # Perform prediction on test data
            y_pred = crf_model.predict(X_test)

            # Evaluate the model
            return str(len(y_pred))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("2")],
        )


@pytest.mark.xfail(
    reason="great_expectations fails often when built from inside a conda environment"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_great_expectations(session, force_install):
    """
    Assert that great_expectations is usable by fetching data context.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["great_expectations"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def ge_test() -> str:
            import great_expectations as gx

            data_context = gx.get_context()
            return str(data_context)

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_simhash(session, force_install):
    """
    Assert that simhash is usable by fetching distance between two simhashes
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["simhash"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def simhash_test() -> str:
            from simhash import Simhash

            return Simhash("aa").distance(Simhash("bb"))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("31")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_eth_abi(session, force_install):
    """
    Assert that eth_abi is usable by encoding two strings to match Solidity's ABI specification.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["eth_abi", "setuptools"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def ethabi_test() -> str:
            from eth_abi import encode

            return str(encode(["bytes32", "bytes32"], [b"a", b"b"]))[:4]

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("b'a\\")],
        )


@pytest.mark.xfail(
    reason="SQL Compilation Error: One or more package conflicts were detected."
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_web3(session, force_install):
    """
    Assert that web3 is usable by .....
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["web3"],
            persist_path=permanent_stage_name,
            force_install=force_install,
            force_push=True,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def web3_test() -> str:
            from web3 import EthereumTesterProvider, Web3

            w3 = Web3(EthereumTesterProvider())
            return str(w3.is_connected())

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("True")],
        )


@pytest.mark.xfail(
    reason="librdkafka C file (required by dependency kafka-confluent) is needed to install hopsworks"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_hopsworks(session, force_install):
    """
    Assert that hopsworks is usable by creating a Hopsworks project object.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["hopsworks"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def hopsworks_test() -> str:
            import hopsworks

            connection = hopsworks.connection()
            project = connection.get_project("my_project")
            return str(project)

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_openai_whisper(session, force_install):
    """
    Assert that openai_whisper is unusable due to native dependencies.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["openai-whisper", "pytorch"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_optuna(session, force_install):
    """
    Assert that optuna is usable by creating a study.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["optuna"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def optuna_test() -> str:
            import optuna

            def objective(trial):
                x = trial.suggest_uniform("x", -10, 10)
                y = trial.suggest_uniform("y", -10, 10)
                return (x - 2) ** 2 + (y + 3) ** 2

            study = optuna.create_study()
            study.optimize(objective, n_trials=100)
            print(str(round(study.best_value, 2)))
            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_ffmpeg(session, force_install):
    """
    Assert that ffmpeg is usable by importing it. (tedious to test further as package requires an audio file)
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["ffmpeg-python"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def ffmpeg_test() -> str:
            import ffmpeg

            print(str(ffmpeg))
            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.xfail(
    reason="Error while dealing with zip files in the pkg_resources: "
    "NotImplementedError: resource_filename() only supported for .egg, not .zip"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_dataprofiler(session, force_install):
    """
    Assert that DataProfiler is usable by profiling simplistic data.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["DataProfiler", "setuptools"],
            persist_path=permanent_stage_name,
            force_install=force_install,
            force_push=True,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def dataprofiler_test() -> str:
            from dataprofiler import Data, Profiler

            data = Data("{'A':'1', 'B':'2'}")
            profile = Profiler(data)
            readable_report = profile.report(
                report_options={"output_format": "compact"}
            )
            return readable_report

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.xfail(reason="ModuleNotFoundError: No module named 'QuantLib'")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_quantlib(session, force_install):
    """
    Assert that quantlib is usable by calculating the price of an option using the Black-Scholes model.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["QuantLib"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def quantlib_test() -> str:
            import QuantLib as ql

            spot_price = 100.0
            strike_price = 100.0
            risk_free_rate = 0.05
            volatility = 0.2
            maturity_date = ql.Date(1, 1, 2024)

            # QuantLib objects
            day_count = ql.Actual365Fixed()
            calendar = ql.UnitedStates()
            calculation_date = calendar.advance(
                ql.Date.todaysDate(), ql.Period(1, ql.Days)
            )
            ql.Settings.instance().evaluationDate = calculation_date

            # Option payoff and exercise
            payoff = ql.PlainVanillaPayoff(ql.Option.Call, strike_price)
            exercise = ql.EuropeanExercise(maturity_date)

            # Option pricing engine
            process = ql.BlackScholesMertonProcess(
                ql.QuoteHandle(ql.SimpleQuote(spot_price)),
                ql.YieldTermStructureHandle(
                    ql.FlatForward(calculation_date, risk_free_rate, day_count)
                ),
                ql.BlackVolTermStructureHandle(
                    ql.BlackConstantVol(
                        calculation_date, calendar, volatility, day_count
                    )
                ),
            )
            option = ql.VanillaOption(payoff, exercise)
            option.setPricingEngine(ql.AnalyticEuropeanEngine(process))

            # Calculate option price
            option_price = option.NPV()
            return str(option_price)

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_routingpy(session, force_install):
    """
    Assert that routingpy is usable by calculating a pedestrian route in Berlin.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["routingpy"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def routingpy_test() -> str:
            import requests
            from routingpy import MapboxValhalla

            # Some locations in Berlin
            coords = [
                [13.413706, 52.490202],
                [13.421838, 52.514105],
                [13.453649, 52.507987],
                [13.401947, 52.543373],
            ]
            client = MapboxValhalla(api_key="mapbox_key")

            try:
                route = client.directions(locations=coords, profile="pedestrian")
                return str(route)
            except requests.exceptions.ConnectionError:
                return "no internet access"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("no internet access")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_httpagentparser(session, force_install):
    """
    Assert that httpagentparser is usable by parsing an agent string.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["httpagentparser"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def httpagentbroswer_test() -> str:
            import httpagentparser

            s = "Mozilla/5.0 (X11; U; Linux i686; en-US) AppleWebKit/532.9 (KHTML, like Gecko) \
                    Chrome/5.0.307.11 Safari/532.9"
            return ",".join(httpagentparser.detect(s).keys())

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("platform,os,bot,browser")],
        )


@pytest.mark.xfail(
    reason="Code error in dependency package 'pulp': AttributeError: 'list' object has no attribute 'split'"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_scikit_criteria(session, force_install):
    """
    Assert that scikit-criteria is usable
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["scikit-criteria"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def scikit_criteria_test() -> str:
            import skcriteria as skc
            from skcriteria.preprocessing import invert_objectives

            matrix = [
                [1, 2, 3],  # alternative 1
                [4, 5, 6],  # alternative 2
            ]

            dm = skc.mkdm(
                matrix,
                [max, max, min],
                weights=[0.5, 0.05, 0.45],
                alternatives=["car 0", "car 1"],
                criteria=["autonomy", "comfort", "price"],
            )
            inverter = invert_objectives.InvertMinimize()
            dmt = inverter.transform(dm)
            return str(dmt["autonomy"])

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.xfail(
    reason="ValueError: The tagger is not opened,  "
    "File 'pycrfsuite/_pycrfsuite.pyx', line 688, in pycrfsuite._pycrfsuite.Tagger.set"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_usaddress(session, force_install):
    """
    Assert that usaddress is usable by parsing a Chicago address.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["usaddress"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def usaddress_test() -> str:
            import usaddress

            return usaddress.parse("123 Main St. Suite 100 Chicago, IL")

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("123")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_nameparser(session, force_install):
    """
    Assert that nameparser is usable by parsing a name.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["nameparser"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def nameparser_test() -> str:
            from nameparser import HumanName

            name = HumanName("Dr. Juan Q. Xavier de la Vega III (Doc Vega)")
            return name.last

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("de la Vega")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_deepchecks(session, force_install):
    """
    Assert that deepchecks package is not usable as it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["deepchecks"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.parametrize("force_install", reinstall_options)
def test_catboost(session, force_install):
    """
    Assert that catboost package is not usable as it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["catboost"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.xfail(reason="Pip error")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_neuralforecast(session, force_install):
    """
    Assert that neuralforecast is usable.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["neuralforecast"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def neuralforecast_test() -> str:
            return "tbd"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("tbd")],
        )


@pytest.mark.xfail(reason="Pip error - needs Rust compiler?")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_allennlp(session, force_install):
    """
    Assert that allennlp is usable.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["allennlp"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def allen_test() -> str:
            return "tbd"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("tbd")],
        )


@pytest.mark.xfail(reason="Requires subprocess")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_pytesseract(session, force_install):
    """
    Assert that pytesseract is usable by performing OCR on an image.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["pytesseract", "Pillow"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def pytesseract_test() -> str:
            import pytesseract
            from PIL import Image, ImageDraw

            # Generate a simple image with text
            image = Image.new("RGB", (200, 100), "white")
            draw = ImageDraw.Draw(image)
            draw.text((50, 40), "Hello!", fill="black")

            # Perform OCR using pytesseract
            text = pytesseract.image_to_string(image)

            # Print the extracted text
            return text

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("Hello!")],
        )


@pytest.mark.xfail(
    reason="RuntimeError: cannot cache function 'sparse_correct_alternative_cosine': no locator available for file..."
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_evidently(session, force_install):
    """
    Assert that evidently is usable by testing data stability for random data.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["evidently"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def evidently_test() -> str:
            import numpy as np
            import pandas as pd
            from evidently.test_preset import DataStabilityTestPreset
            from evidently.test_suite import TestSuite

            # Generate random data
            np.random.seed(42)
            random_data = np.random.rand(120, 4)
            column_names = ["feature1", "feature2", "feature3", "feature4"]
            random_dataframe = pd.DataFrame(data=random_data, columns=column_names)

            # Split the random data into current and reference data
            current_data = random_dataframe.iloc[:60]
            reference_data = random_dataframe.iloc[60:]

            # Create the test suite and run tests
            data_stability = TestSuite(tests=[DataStabilityTestPreset()])
            data_stability.run(
                current_data=current_data,
                reference_data=reference_data,
                column_mapping=None,
            )

            return data_stability.json()

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_ua_parser(session, force_install):
    """
    Assert that ua-parser is usable by deducing web browser name from a user agent string.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["ua-parser"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def uaparser_test() -> str:
            from ua_parser import user_agent_parser

            ua_string = "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_9_4) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/41.0.2272.104 Safari/537.36"
            data = user_agent_parser.Parse(ua_string)
            return data["user_agent"]["family"]

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("Chrome")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_azure_keyvault(session, force_install):
    """
    Assert that azure keyvault is usable.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["azure-keyvault-secrets", "azure-identity"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def azure_kv_test() -> str:
            from azure.identity import DefaultAzureCredential
            from azure.keyvault.secrets import SecretClient

            credential = DefaultAzureCredential()
            secret_client = SecretClient(
                vault_url="https://my-key-vault.vault.azure.net/", credential=credential
            )
            print(secret_client)
            return "worked"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("worked")],
        )


def test_pandas_dedupe(session):
    """
    Assert that pandas_dedupe package is not usable as it contains native code.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        with pytest.raises(RuntimeError) as ex_info:
            session.add_packages(
                ["pandas_dedupe"],
                persist_path=permanent_stage_name,
                force_install=True,
            )
        assert "Your code depends on native dependencies" in str(ex_info)


@pytest.mark.xfail(
    reason='StopIteration: _available_dir = [p for p in next(os.walk(_module_path))[1] if not p.startswith("__]")]'
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_moving_pandas(session, force_install):
    """
    Assert that moving_pandas package is usable.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["movingpandas", "geopandas"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def movingpandas_test() -> str:
            from datetime import datetime, timedelta

            import pandas as pd
            from movingpandas.trajectory import Trajectory
            from shapely.geometry import Point

            # Generate some sample data
            data = {"timestamp": [], "latitude": [], "longitude": []}
            start_time = datetime(2023, 1, 1)
            num_points = 100
            interval = timedelta(minutes=10)

            for i in range(num_points):
                timestamp = start_time + i * interval
                latitude = 40.0 + i * 0.01
                longitude = -90.0 + i * 0.01
                data["timestamp"].append(timestamp)
                data["latitude"].append(latitude)
                data["longitude"].append(longitude)

            df = pd.DataFrame(data)
            df["geometry"] = [Point(xy) for xy in zip(df["longitude"], df["latitude"])]

            # Create a Trajectory object
            trajectory = Trajectory(df, "timestamp", "geometry")

            # Test some functionalities of the Trajectory object
            print(f"Number of points in the trajectory: {len(trajectory)}")
            print(f"Start time of the trajectory: {trajectory.get_start_time()}")
            print(f"End time of the trajectory: {trajectory.get_end_time()}")
            print(f"Total duration of the trajectory: {trajectory.get_duration()}")
            print(f"Bounding box of the trajectory: {trajectory.get_bbox()}")
            print(f"Spatial extent of the trajectory: {trajectory.get_extent()}")
            print(f"Speeds of the trajectory: {trajectory.get_speed()}")

            return "tbd"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("tbd")],
        )


@pytest.mark.xfail(reason="Pip error, package not found")
@pytest.mark.parametrize("force_install", reinstall_options)
def test_cplex(session, force_install):
    """
    Assert that cplex is usable
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["cplex"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def cplex_test() -> str:
            import cplex

            # Coefficients for the objective function
            objective_coefficients = [1.0, 2.0, 3.0]

            # Constraint coefficients
            constraint_coefficients = [[1.0, 1.0, 1.0], [2.0, -1.0, 1.0]]

            # Constraint senses ('L' for less than or equal to, 'G' for greater than or equal to, 'E' for equal to)
            constraint_senses = ["L", "L"]

            # Right-hand side values for the constraints
            rhs_values = [5.0, 3.0]

            # Variable bounds
            lower_bounds = [0.0, 0.0, 0.0]
            upper_bounds = [cplex.infinity, cplex.infinity, cplex.infinity]

            # Create an instance of Cplex
            problem = cplex.Cplex()

            # Set the objective sense to maximize
            problem.objective.set_sense(problem.objective.sense.maximize)

            # Add variables to the problem
            problem.variables.add(
                obj=objective_coefficients, lb=lower_bounds, ub=upper_bounds
            )

            # Add constraints to the problem
            problem.linear_constraints.add(
                lin_expr=constraint_coefficients,
                senses=constraint_senses,
                rhs=rhs_values,
            )

            # Solve the problem
            problem.solve()

            # Get the solution status
            solution_status = problem.solution.get_status_string()

            # Get the objective value
            objective_value = problem.solution.get_objective_value()

            # Get the variable values
            variable_values = problem.solution.get_values()

            # Print the results
            print("Solution Status: ", solution_status)
            print("Objective Value: ", objective_value)
            print("Variable Values: ", variable_values)

            return str((solution_status, objective_value, variable_values))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("tbd")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_forex_python(session, force_install):
    """
    Assert that forex-python is importable (not usable as it requires an internet connection).
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["forex-python"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def forex_test() -> str:
            import requests
            from forex_python.converter import CurrencyRates

            c = CurrencyRates()
            try:
                return c.get_rates("USD")
            except requests.exceptions.ConnectionError:
                return "no internet"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("no internet")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_chainladder(session, force_install):
    """
    Assert that chainladder is usable by calculating cumulative triangle.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["chainladder"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def chainladder_test() -> str:
            import pandas as pd
            from chainladder import Triangle

            # Create the data dictionary with the desired columns
            data_dict = {
                "Acc Year": [2018, 2018, 2019, 2019],
                "Cal Year": [2018, 2018, 2019, 2019],
                "Cal Month": [1, 2, 1, 2],
                "Paid Loss": [1000, 2000, 1500, 2500],
            }

            # Create the DataFrame
            data = pd.DataFrame(data_dict)
            triangle = Triangle(
                data,
                origin="Acc Year",
                development=["Cal Year", "Cal Month"],
                columns=["Paid Loss"],
            )
            return str(len(triangle.to_frame()))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("2")],
        )


@pytest.mark.xfail(
    reason="Please setup your environment variables gcc/gxx with your GCC/CLANG path"
)
@pytest.mark.parametrize("force_install", reinstall_options)
def test_pyfhel(session, force_install):
    """
    Assert that pyfhel is usable by encrypting two integers.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["Pyfhel"],
            persist_path=permanent_stage_name,
            force_install=force_install,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def pyfhel_test() -> str:
            import numpy as np
            from Pyfhel import Pyfhel

            HE = Pyfhel()  # Creating empty Pyfhel object
            HE.contextGen(
                scheme="bfv", n=2**14, t_bits=20
            )  # Generate context for 'bfv'/'ckks' scheme
            HE.keyGen()  # Key Generation: generates a pair of public/secret keys

            integer1 = np.array([127], dtype=np.int64)
            integer2 = np.array([-2], dtype=np.int64)
            ctxt1 = HE.encryptInt(integer1)  # Encryption makes use of the public key
            ctxt2 = HE.encryptInt(
                integer2
            )  # For integers, encryptInt function is used.

            ctxtSum = ctxt1 + ctxt2
            ctxtSub = ctxt1 - ctxt2
            ctxtMul = ctxt1 * ctxt2

            resSum = HE.decryptInt(
                ctxtSum
            )  # Decryption must use the corresponding function
            #  decryptInt.
            resSub = HE.decryptInt(ctxtSub)
            resMul = HE.decryptInt(ctxtMul)

            return str(len(resSum) + len(resSub) + len(resMul))

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("TBD")],
        )


@pytest.mark.parametrize("force_install", reinstall_options)
def test_sqlalchemy(session, force_install):
    """
    Assert that snowflake-sqlalchemy is usable.
    """
    with patch.object(session, "_is_anaconda_terms_acknowledged", lambda: True):
        session.add_packages(
            ["snowflake-sqlalchemy"],
            persist_path=permanent_stage_name,
            force_install=force_install,
            force_push=True,
        )
        udf_name = Utils.random_name_for_temp_object(TempObjectType.FUNCTION)

        @udf(name=udf_name, session=session)
        def sqlalchemy_test() -> str:
            from sqlalchemy import create_engine

            print(str(create_engine))
            return "works"

        Utils.check_answer(
            session.sql(f"select {udf_name}()").collect(),
            [Row("works")],
        )
